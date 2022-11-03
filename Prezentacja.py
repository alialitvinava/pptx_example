import pptx
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import ColorFormat, RGBColor
import datetime

import pandas as pd
import numpy as np

class Prezentacja:

    def __init__(self, kampania):
        self.kampania = str(kampania)  #atrybut klasy Prezentacja
        
    def ładowanie_danych(self):            #ładowanie danych w postaci data frame w pandas
        df = pd.read_excel('dane_sgh.xlsx')
        data = df.fillna(0)                #wypełnienie pustych komórek data frame zerami
        return data                        #zwrot data frame

    def obróbka_danych(self, data, kampania):  #przetworzenie danych: filtrowanie wybranej kampanii z data frame i wygenerowanie tablic przestawnych dla wykresów
        kampania = self.kampania
        wydarzenie = data.loc[data['Kampania'] == kampania] # filtrowanie
        przestawna_miejsca = wydarzenie.pivot_table(values='Aktywni uczestnicy', index='Miejsce', #utworzenie tablicy przestawnej z miejscami kampanii i liczbą aktywnych uczestników
                                                    aggfunc=np.sum).reset_index()
        przestawna_miejsca_koszty = wydarzenie.pivot_table(values=['Aktywni uczestnicy', 'Koszty'], index='Miejsce',
                                                           aggfunc=np.sum).reset_index()            #utworzenie tablicy przestawnej z miejscami kampanii, kosztami i liczbą aktywnych ucze
        przestawna_miejsca_koszty['przeciętne_koszty'] = przestawna_miejsca_koszty['Koszty'] / przestawna_miejsca_koszty['Aktywni uczestnicy'] # dodawanie nowej zmiennej na podstawie obliczeń z kolumnami
        przestawna_miejsca_organizatorzy = wydarzenie.pivot_table(values=['Liczba organizatorów'], index='Miejsce',
                                                                  aggfunc=np.sum).reset_index()  #utworzenie tablicy przestawnej z miejscami kampanii i liczbą organizatorów
        atrybuty_pivot = wydarzenie.pivot_table(values=['Aktywni uczestnicy', 'Liczba osób w miejscu'],
                                                index=['Atrybuty'], aggfunc=np.sum).reset_index()  #utworzenie tablicy przestawnej z atrybutami kampanii i liczbą osób w miejscu i aktywnymi uczestnikami
        atrybuty_pivot['udział_aktywnych'] = atrybuty_pivot['Aktywni uczestnicy'] / atrybuty_pivot['Liczba osób w miejscu']   # dodawanie nowej kolumny z obliczeniem udziału aktywnych uczestników wśród całej liczby osob w miejscu w tym dniu
        atrybuty_miejsca_pivot = wydarzenie.pivot_table(values=['Aktywni uczestnicy', 'Liczba osób w miejscu'],
                                                        index=['Atrybuty', 'Miejsce'], aggfunc=np.sum).reset_index() #utworzenie tablicy przestawnej z atrybutami kampanii, miejscem, liczbą osób w miejscu i aktywnymi uczestnikami
        atrybuty_miejsca_pivot['udział_aktywnych'] = atrybuty_miejsca_pivot['Aktywni uczestnicy'] / atrybuty_miejsca_pivot['Liczba osób w miejscu']  # dodawanie nowej kolumny z obliczeniem udziału aktywnych uczestników wśród całej liczby osob w miejscu w tym dniu
        miejsca_słownik = {'Budynek G': "budynku G", 'Budynek A': "budynku A",
                           'Budynek Sabinki': "budynku Sabinki", 'Budynek M': "budynku M",
                           'Budynek Grosik': "budynku Grosik", 'Budynek W': "budynku W", 'Budynek C': "budynku C",
                           'Biblioteka': "bibliotece"} # słownik z odmianą nazw miejsc
        przestawna_miejsca['index'] = przestawna_miejsca['Miejsce'].map(miejsca_słownik) #łączenie słownika z tabelą przestawną
        przestawna_miejsca.set_index('index', inplace=True)
        przestawna_miejsca_koszty['index'] = przestawna_miejsca_koszty['Miejsce'].map(miejsca_słownik) #łączenie słownika z tabelą przestawną
        przestawna_miejsca_koszty.set_index('index', inplace=True)
        przestawna_miejsca_organizatorzy['index'] = przestawna_miejsca_organizatorzy['Miejsce'].map(miejsca_słownik) #łączenie słownika z tabelą przestawną
        przestawna_miejsca_organizatorzy.set_index('index', inplace=True)
        return przestawna_miejsca, przestawna_miejsca_koszty, przestawna_miejsca_organizatorzy, atrybuty_pivot, atrybuty_miejsca_pivot

    def utwórz_slide(self, presentation, title, layout=4):  #funkcja, która bierze ze wzorca prezentacji slajd o numerze 4 i wkleja i nadaje placeholderowi z tytułem kształt tytułu + wkleja tekst do tego placeholderu
       layout = presentation.slide_layouts[4]
       slide = presentation.slides.add_slide(layout)
       if title is not None:
           slide_title = slide.shapes.title
           slide_title.text = title
       return slide


    def utwórz_pptx(self, przestawna_miejsca, przestawna_miejsca_koszty, przestawna_miejsca_organizatorzy, atrybuty_pivot,
                     atrybuty_miejsca_pivot, kampania):
        presentation = pptx.Presentation("SGH_prezentacja_ogolna_PRZYKLAD_pl.pptx") #ładowanie wzorca prezentacji
        title_slide_layout = presentation.slide_layouts[1]               # ze wzorca slajdu jest wybierany 1 slajd jako slajd tytułowu
        slide = presentation.slides.add_slide(title_slide_layout)
        title = slide.shapes.title                                       #definiowanie placeholderu tytułu
        title.text = kampania                                            #wklejenie do placeholderu z tytułem text z nazwą kampanii
        subtitle = slide.placeholders[1]                                  #definiowanie placeholderu podtytułu
        date = datetime.datetime.today().strftime('%Y-%m-%d')            #dzisiejsza data
        subtitle.text = f"{date}"                                          #wklejenie do placeholderu z podtytułem dzisiejszą datę jako tekst

        slide1 = self.utwórz_slide(presentation, "Liczba Aktywnych Użytkowników w każdym budynku") #tworzenie slajdów z określonym tytułem oraz wklejenie do niego wykresu
        self.create_chart_slide_1(przestawna_miejsca, slide1)

        slide2 = self.utwórz_slide(presentation, "Liczba użytkowników i koszty wydarzenia dla każdego budynku") #tworzenie slajdów z określonym tytułem oraz wklejenie do niego wykresu
        self.create_chart_slide_2(przestawna_miejsca_koszty, slide2)

        slide3 = self.utwórz_slide(presentation, "Przeciętne koszty za aktywnego użytkownika") #tworzenie slajdów z określonym tytułem oraz wklejenie do niego wykresu
        self.create_chart_slide_3(przestawna_miejsca_koszty, slide3)

        slide4 = self.utwórz_slide(presentation, "Liczba Organizatorów w poszczególnych budynkach") #tworzenie slajdów z określonym tytułem oraz wklejenie do niego wykresu
        self.create_chart_slide_4(przestawna_miejsca_organizatorzy, slide4)

        slide5 = self.utwórz_slide(presentation, "Udział Aktywnych użytkowników dla poszczególnych typów kampanii oraz dla poszczególnych budynków") #tworzenie slajdów z określonym tytułem oraz wklejenie do niego wykresu
        self.create_chart_slide_5(atrybuty_pivot, atrybuty_miejsca_pivot, slide5)

        presentation.save(f'{kampania}_{date}.pptx')  #pobieranie gotowej prezentacji

    def create_chart_slide_1(self, przestawna_miejsca, slide): #tworzenie wykresu i tekstu do slajdu 1
        chart_data_1 = CategoryChartData()
        chart_data_1.categories = list(przestawna_miejsca['Miejsce'])
        chart_data_1.add_series('Liczba Aktywnych Użytkowników', list(przestawna_miejsca['Aktywni uczestnicy']))
        CHART_TYPE = XL_CHART_TYPE.COLUMN_CLUSTERED
        chart_left = Inches(0.57)
        chart_top = Inches(1.36)
        chart_width = Inches(5.94)
        chart_height = Inches(3.67)
        chart_1 = slide.shapes.add_chart(CHART_TYPE, chart_left, chart_top, chart_width, chart_height,
                                         chart_data_1).chart
        plot_1 = chart_1.plots[0]
        plot_1.has_data_labels = True
        data_labels = plot_1.data_labels
        data_labels.font.size = Pt(7)
        value_axis = chart_1.value_axis
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(11)
        category_axis = chart_1.category_axis
        category_axis.tick_labels.font.size = Pt(14)
        chart_1.chart_title.has_text_frame = True
        chart_1.chart_title.text_frame.text = 'Liczba Uczestników'

        maxValueIndex = przestawna_miejsca['Aktywni uczestnicy'].idxmax()
        left = Inches(6.89)
        top = Inches(2.03)
        width = Inches(2.68)
        height = Inches(0.87)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Najwięcej aktywnych uczestników przybywało w {maxValueIndex}."
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        line = txBox.line
        line.color.rgb = RGBColor(41, 253, 66)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p.font.size = Pt(11)

    def create_chart_slide_2(self, przestawna_miejsca_koszty, slide):  #tworzenie wykresu i tekstu do slajdu 2
        chart_data = CategoryChartData()
        chart_data.categories = list(przestawna_miejsca_koszty['Miejsce'])
        chart_data.add_series('Koszty', list(przestawna_miejsca_koszty['Koszty'].round(2)))
        CHART_TYPE = XL_CHART_TYPE.COLUMN_CLUSTERED
        chart_left = Inches(0.57)
        chart_top = Inches(1.36)
        chart_width = Inches(5.94)
        chart_height = Inches(3.67)
        chart_1 = slide.shapes.add_chart(CHART_TYPE, chart_left, chart_top, chart_width, chart_height, chart_data).chart
        plot_1 = chart_1.plots[0]
        plot_1.has_data_labels = True
        data_labels = plot_1.data_labels
        data_labels.number_format = "0\ zł"
        data_labels.font.size = Pt(7)
        value_axis = chart_1.value_axis
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(11)
        category_axis = chart_1.category_axis
        category_axis.tick_labels.font.size = Pt(14)
        chart_1.chart_title.has_text_frame = True
        chart_1.chart_title.text_frame.text = 'Koszty kampanii w poszczególnych budynkach'

        maxValueIndex = przestawna_miejsca_koszty['Koszty'].idxmax()
        left = Inches(6.89)
        top = Inches(2.03)
        width = Inches(2.68)
        height = Inches(0.87)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Najwyższe koszty osiągnięte są w {maxValueIndex}."
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        line = txBox.line
        line.color.rgb = RGBColor(41, 253, 66)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p.font.size = Pt(11)

    def create_chart_slide_3(self, przestawna_miejsca_koszty, slide):  #tworzenie wykresu i tekstu do slajdu 3
        chart_data = CategoryChartData()
        chart_data.categories = list(przestawna_miejsca_koszty['Miejsce'])
        chart_data.add_series('Przeciętne koszty', list(przestawna_miejsca_koszty['przeciętne_koszty'].round(2)))
        CHART_TYPE = XL_CHART_TYPE.COLUMN_CLUSTERED
        chart_left = Inches(0.57)
        chart_top = Inches(1.36)
        chart_width = Inches(5.94)
        chart_height = Inches(3.67)
        chart_1 = slide.shapes.add_chart(CHART_TYPE, chart_left, chart_top, chart_width, chart_height, chart_data).chart
        plot_1 = chart_1.plots[0]
        plot_1.vary_by_categories = False
        plot_1.has_data_labels = True
        data_labels = plot_1.data_labels
        data_labels.number_format = "0\ zł"
        data_labels.font.size = Pt(7)
        value_axis = chart_1.value_axis
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(11)
        category_axis = chart_1.category_axis
        category_axis.tick_labels.font.size = Pt(14)
        chart_1.chart_title.has_text_frame = True
        chart_1.chart_title.text_frame.text = 'Przeciętne koszty za uczestnika'

        maxValueIndex = przestawna_miejsca_koszty['przeciętne_koszty'].idxmax()
        left = Inches(6.89)
        top = Inches(2.03)
        width = Inches(2.68)
        height = Inches(0.87)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"Najwyższe przeciętne koszty za uczestnika osiągnięte są w {maxValueIndex}."
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        line = txBox.line
        line.color.rgb = RGBColor(41, 253, 66)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p.font.size = Pt(11)

    def create_chart_slide_4(self, przestawna_miejsca_organizatorzy, slide):  #tworzenie wykresu i tekstu do slajdu 4
        chart_data_2 = ChartData()
        chart_data_2.categories = list(przestawna_miejsca_organizatorzy['Miejsce'])
        chart_data_2.add_series('Liczba Organizatorów', list(przestawna_miejsca_organizatorzy['Liczba organizatorów']))
        CHART_TYPE = XL_CHART_TYPE.PIE
        chart_left = Inches(0.57)
        chart_top = Inches(1.36)
        chart_width = Inches(5.94)
        chart_height = Inches(3.67)
        chart_2 = slide.shapes.add_chart(CHART_TYPE, chart_left, chart_top, chart_width, chart_height,
                                         chart_data_2).chart
        plot_2 = chart_2.plots[0]
        plot_2.has_data_labels = True
        data_labels = plot_2.data_labels
        data_labels.font.size = Pt(7)
        chart_2.chart_title.has_text_frame = True
        chart_2.chart_title.text_frame.text = 'Liczba organizatorów w poszczególnych budynkach'
        chart_2.has_legend = True
        chart_2.legend.position = XL_LEGEND_POSITION.RIGHT
        chart_2.legend.include_in_layout = False
        chart_2.legend.font.size = Pt(11)

        maxValueIndex = przestawna_miejsca_organizatorzy['Liczba organizatorów'].idxmax()
        left = Inches(6.89)
        top = Inches(2.03)
        width = Inches(2.68)
        height = Inches(0.87)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Najwięcej organizatorów było w {maxValueIndex}."
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        line = txBox.line
        line.color.rgb = RGBColor(41, 253, 66)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p.font.size = Pt(11)

    def create_chart_slide_5(self, atrybuty_pivot, atrybuty_miejsca_pivot, slide):  #tworzenie wykresów i tekstu do slajdu 5
        chart_data_1 = CategoryChartData()
        chart_data_1.categories = list(atrybuty_pivot['Atrybuty'])
        chart_data_1.add_series('Udział aktywnych użytkowników', list(atrybuty_pivot['udział_aktywnych'].round(1)))
        CHART_TYPE = XL_CHART_TYPE.COLUMN_CLUSTERED
        chart_left = Inches(0.26)
        chart_top = Inches(1.74)
        chart_width = Inches(4.67)
        chart_height = Inches(2.74)
        chart_1 = slide.shapes.add_chart(CHART_TYPE, chart_left, chart_top, chart_width, chart_height,
                                         chart_data_1).chart
        plot_1 = chart_1.plots[0]
        plot_1.vary_by_categories = False
        plot_1.has_data_labels = True
        data_labels = plot_1.data_labels
        data_labels.number_format = "0.0\%"
        data_labels.font.size = Pt(6)
        value_axis = chart_1.value_axis
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(9)
        category_axis = chart_1.category_axis
        category_axis.tick_labels.font.size = Pt(9)
        chart_1.chart_title.has_text_frame = True
        chart_1.chart_title.text_frame.text = 'Udział aktywnych użytkowników'

        df_chart_5 = atrybuty_miejsca_pivot.pivot_table(values="udział_aktywnych", index="Atrybuty", columns="Miejsce")
        df_chart_5.reset_index(inplace=True)
        df_chart_6 = df_chart_5.fillna(0)
        chart_data_2 = CategoryChartData()
        chart_data_2.categories = list(df_chart_6['Atrybuty'])
        columns = df_chart_6.columns.values.tolist()
        for c in columns:
            series_title = "%s" % c
            series_values = list(df_chart_6["%s" % c])
            chart_data_2.add_series(series_title, series_values)
        CHART_TYPE = XL_CHART_TYPE.COLUMN_CLUSTERED
        chart_left_2 = Inches(4.87)
        chart_top_2 = Inches(1.74)
        chart_width_2 = Inches(4.67)
        chart_height_2 = Inches(2.74)
        chart_2 = slide.shapes.add_chart(CHART_TYPE, chart_left_2, chart_top_2, chart_width_2, chart_height_2,
                                         chart_data_2).chart
        plot_2 = chart_2.plots[0]
        plot_2.has_data_labels = True
        data_labels = plot_2.data_labels
        data_labels.number_format = "0.0\%"
        data_labels.font.size = Pt(6)
        value_axis = chart_2.value_axis
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(9)
        category_axis = chart_2.category_axis
        category_axis.tick_labels.font.size = Pt(9)
        chart_2.chart_title.has_text_frame = True
        chart_2.chart_title.text_frame.text = 'Udział Aktywnych Uczestników'
        chart_2.has_legend = True
        chart_2.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart_2.legend.include_in_layout = False
        chart_2.legend.font.size = Pt(6)

        atrybuty_pivot.set_index('Atrybuty', inplace=True)
        maxValueIndex = atrybuty_pivot['udział_aktywnych'].idxmax()
        left = Inches(2.36)
        top = Inches(4.48)
        width = Inches(4.7)
        height = Inches(0.59)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Najbardziej atrakcyjnymi dla uczestników były {maxValueIndex}."
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        line = txBox.line
        line.color.rgb = RGBColor(41, 253, 66)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p.font.size = Pt(11)

    def main(self):       #wywoływanie tworzenia prezentacji
        kampania = self.kampania

        data = self.ładowanie_danych()

        przestawna_miejsca, przestawna_miejsca_koszty, przestawna_miejsca_organizatorzy, atrybuty_pivot, atrybuty_miejsca_pivot = self.obróbka_danych(data, kampania)

        self.utwórz_pptx(przestawna_miejsca, przestawna_miejsca_koszty, przestawna_miejsca_organizatorzy, atrybuty_pivot,
                     atrybuty_miejsca_pivot, kampania)
